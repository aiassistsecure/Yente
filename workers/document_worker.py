#!/usr/bin/env python3
"""Bytes in, text out. Nothing else.

WHY THIS IS A SEPARATE PROCESS AND NOT A LIBRARY CALL

Two reasons, and the second is the one that matters.

1. NEDB's data directory is locked exclusively and the engine exposes no
   `close()`. A second process holding it is a split-brain by construction. So a
   worker that touched the store would have to be the daemon, and then a
   segfault in a PDF parser takes the mailbox down with it.

2. An attachment is an untrusted file from a stranger, and document parsers are
   the largest attack surface in this entire system — PDF and Office parsers have
   long CVE histories, and the interesting ones are memory-safety bugs, not logic
   bugs. A crash here should cost one subprocess and one attachment.

THIS PROCESS CAN DO NOTHING

It has no network client, no store handle, no credentials, and no way to reach
the graph. It reads bytes on stdin, writes JSON on stdout, exits. That is the
whole contract. §5's "attachment workers must not execute macros, scripts,
binaries, embedded executables, or instructions found inside documents" is
enforced by there being nothing here to execute them WITH.

PARSER CHOICE IS A SECURITY DECISION

Every library below was picked for what it refuses to do:

  pdfminer.six  - pure Python text extraction. No JavaScript engine, no form
                  actions, no external resource fetching. A PDF's /OpenAction
                  is data it will never run.
  python-docx   - walks document XML. Ignores embedded objects and VBA rather
                  than resolving them; a macro in a .docm is inert bytes.
  openpyxl      - read_only=True AND data_only=True. The second matters more
                  than it looks: it reads cached values instead of evaluating
                  formulas, so a sheet full of WEBSERVICE() or DDE calls yields
                  text, not requests.
  python-pptx   - same XML-walking approach as docx.

No OCR. §4 says prefer deterministic native extraction, and an OCR pass is a
second, much larger dependency that produces lower-quality text with no
provenance. A scanned PDF returns empty text and says why, which is a true and
useful answer.

PROVENANCE IS PART OF THE OUTPUT

Page, sheet and slide markers are emitted inline as `[[page 3]]`. That is not
decoration: a claim quoting an attachment needs to say WHERE, and §3's example
is literally "Attachment - fundraising.pdf - page 3". Because the markers live
in the same text the span verifier checks against, a quote can include one and
still ground.

USAGE
    echo '{"filename":"x.pdf","mime_type":"application/pdf","b64":"..."}' \\
      | python3 workers/document_worker.py

Always exits 0 with a JSON verdict. A parser that raised is a fact about the
document, not a reason to make the caller handle a signal.
"""
from __future__ import annotations

import base64
import io
import json
import os
import sys

# A ceiling, because an attachment is attacker-controlled and a zip bomb is a
# real thing. 25MB of business document is already unusual.
def _int_env(name: str, default: int) -> int:
    """An UNSET var and an EMPTY var are both "use the default".

    `int(os.environ.get(name, default))` raises on an empty string, and a caller
    that passes `VAR: ""` to mean "not set" is a completely reasonable caller —
    ours did. The traceback was also invisible until the Node side stopped
    parsing empty stdout as an empty verdict, so this failed twice over.
    """
    raw = (os.environ.get(name) or "").strip()
    try:
        return int(raw) if raw else default
    except ValueError:
        return default


MAX_BYTES = _int_env("YENTE_DOC_MAX_BYTES", 25 * 1024 * 1024)

# Per-extractor output cap. A 4,000-page PDF should not become a prompt.
MAX_CHARS = _int_env("YENTE_DOC_MAX_CHARS", 200_000)


def _truncate(text: str) -> tuple[str, bool]:
    if len(text) <= MAX_CHARS:
        return text, False
    return text[:MAX_CHARS], True


# --- extractors ------------------------------------------------------------
# Each returns (text, structure). `structure` is what the caller can show a
# person: how many pages, which sheets. Imports are function-local so a box
# missing one library still handles every other type.


def extract_pdf(data: bytes) -> tuple[str, dict]:
    from pdfminer.high_level import extract_text_to_fp
    from pdfminer.layout import LAParams

    out = io.StringIO()
    # No password, no JS, no external resources — extract_text_to_fp has no code
    # path to any of them.
    extract_text_to_fp(io.BytesIO(data), out, laparams=LAParams(), output_type="text")
    raw = out.getvalue()

    # pdfminer separates pages with \f. Turn that into a marker a quote can
    # legitimately contain, so "page 3" is checkable rather than asserted.
    pages = raw.split("\f")
    parts = []
    for index, page in enumerate(pages, start=1):
        page = page.strip()
        if page:
            parts.append(f"[[page {index}]]\n{page}")
    return "\n\n".join(parts), {"pages": len([p for p in pages if p.strip()])}


def extract_docx(data: bytes) -> tuple[str, dict]:
    import docx

    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]

    # Tables carry the facts in a lot of business documents — a rate card, a
    # team list — and dropping them silently would lose exactly the structured
    # content §4 asks us to keep.
    tables = 0
    for table in document.tables:
        tables += 1
        parts.append(f"[[table {tables}]]")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    return "\n".join(parts), {"paragraphs": len(document.paragraphs), "tables": tables}


def extract_xlsx(data: bytes) -> tuple[str, dict]:
    import openpyxl

    # read_only: never load the whole workbook graph.
    # data_only: read CACHED VALUES, never evaluate formulas. A sheet full of
    # WEBSERVICE() or DDE becomes text instead of outbound requests.
    book = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    sheets = []
    for sheet in book.worksheets:
        sheets.append(sheet.title)
        parts.append(f"[[sheet {sheet.title}]]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    book.close()
    return "\n".join(parts), {"sheets": sheets}


def extract_pptx(data: bytes) -> tuple[str, dict]:
    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    parts = []
    for index, slide in enumerate(deck.slides, start=1):
        parts.append(f"[[slide {index}]]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts), {"slides": len(deck.slides)}


def extract_text(data: bytes) -> tuple[str, dict]:
    # errors="replace": a mislabelled encoding is not a reason to lose a
    # message. A replacement character in one word beats discarding the file.
    return data.decode("utf-8", errors="replace"), {}


def extract_csv(data: bytes) -> tuple[str, dict]:
    import csv

    text = data.decode("utf-8", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(text[:4096])
    except Exception:
        dialect = csv.excel
    rows = list(csv.reader(io.StringIO(text), dialect))
    return "\n".join(" | ".join(r) for r in rows if any(r)), {"rows": len(rows)}


EXTRACTORS = {
    "application/pdf": extract_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
    "text/plain": extract_text,
    "text/markdown": extract_text,
    "text/csv": extract_csv,
    "text/tab-separated-values": extract_csv,
    "application/json": extract_text,
}

# Extension fallback: senders and mail clients mislabel MIME types constantly,
# and refusing a real .docx because Outlook called it octet-stream would be a
# worse product for no security gain — the extractor is chosen by us either way.
BY_EXTENSION = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".tsv": "text/tab-separated-values",
    ".json": "application/json",
}

# Refused BY NAME rather than falling through to "unsupported", so the log says
# something true about what arrived. These are executables and archives: there is
# no text to extract and every one of them is a delivery mechanism.
DANGEROUS = {
    ".exe", ".dll", ".so", ".dylib", ".bat", ".cmd", ".com", ".scr", ".ps1",
    ".vbs", ".js", ".jar", ".msi", ".app", ".deb", ".rpm", ".sh",
    ".docm", ".xlsm", ".pptm",   # macro-enabled Office
    ".zip", ".rar", ".7z", ".tar", ".gz",   # containers: unpacking is a decision, not a default
}


def resolve(filename: str, mime_type: str):
    lowered = (filename or "").lower()
    ext = os.path.splitext(lowered)[1]

    if ext in DANGEROUS:
        return None, f"refused by type ({ext}): executable, macro-enabled or archive"

    if mime_type in EXTRACTORS:
        return EXTRACTORS[mime_type], None
    if ext in BY_EXTENSION:
        return EXTRACTORS[BY_EXTENSION[ext]], None
    return None, f"unsupported type: {mime_type or ext or 'unknown'}"


def main() -> int:
    try:
        request = json.loads(sys.stdin.read() or "{}")
    except Exception as exc:
        print(json.dumps({"ok": False, "error": f"bad request: {exc}"}))
        return 0

    filename = request.get("filename") or ""
    mime_type = request.get("mime_type") or ""

    try:
        data = base64.b64decode(request.get("b64") or "", validate=False)
    except Exception as exc:
        print(json.dumps({"ok": False, "filename": filename, "error": f"bad base64: {exc}"}))
        return 0

    if not data:
        print(json.dumps({"ok": False, "filename": filename, "error": "empty attachment"}))
        return 0

    if len(data) > MAX_BYTES:
        print(json.dumps({
            "ok": False, "filename": filename,
            "error": f"too large: {len(data)} bytes exceeds {MAX_BYTES}",
        }))
        return 0

    extractor, refusal = resolve(filename, mime_type)
    if refusal:
        print(json.dumps({"ok": False, "filename": filename, "error": refusal}))
        return 0

    try:
        text, structure = extractor(data)
    except Exception as exc:
        # A parser that raised is a FACT ABOUT THE DOCUMENT, reported as data.
        # Crashing would make the caller handle a signal for something that is
        # simply an unreadable file.
        print(json.dumps({
            "ok": False, "filename": filename,
            "error": f"{type(exc).__name__}: {exc}"[:400],
        }))
        return 0

    text, truncated = _truncate(text.strip())
    print(json.dumps({
        "ok": True,
        "filename": filename,
        "mime_type": mime_type,
        "bytes": len(data),
        "text": text,
        "chars": len(text),
        "truncated": truncated,
        "structure": structure,
        # An empty result is a real answer, not a failure. Almost always a
        # scanned document — and saying so is more useful than a silent nothing,
        # because it tells you OCR is the missing piece rather than the parser.
        "empty_reason": None if text else "no extractable text (likely a scan; OCR is not enabled)",
    }))
    return 0


if __name__ == "__main__":
    sys.exit(main())
